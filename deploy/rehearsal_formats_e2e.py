"""Destructive format rehearsal for an isolated LectureSift migration database.

This script is intentionally not imported by the service.  It creates four
small, synthetic jobs against the local API, proves that Celery and private
object storage complete the durable path, and then removes the rehearsal
account and every object created by the run.
"""

from __future__ import annotations

import atexit
import json
import os
from pathlib import Path
from pathlib import PurePosixPath, PureWindowsPath
import re
import secrets
import shutil
import stat
import subprocess
import tempfile
import time
from typing import Any
import uuid
import zipfile

import httpx
from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from sqlalchemy.engine import make_url

from lecturesift import config
from lecturesift.billing_service import register_user, verify_email
from lecturesift.jobs import JOBS
from lecturesift.rollout_service import admin_close_user_account, admin_set_user_subscription
from lecturesift.storage import STORAGE


BASE_URL = "http://127.0.0.1:8000"
TERMINAL_TIMEOUT_SECONDS = max(
    300,
    min(3600, int(os.getenv("LECTURESIFT_REHEARSAL_FORMAT_TIMEOUT_SECONDS", "1500"))),
)
CURRENT_STAGE = "startup"


class FormatRehearsalFailure(RuntimeError):
    """A deliberately small, public-safe failure from one synthetic case."""

    def __init__(self, case_name: str, public_error_code: str) -> None:
        allowed_cases = {"native_documents", "ocr_images", "mp3_audio", "mp4_video"}
        safe_case = case_name if case_name in allowed_cases else "unknown"
        safe_code = (
            public_error_code
            if re.fullmatch(r"LS-[A-Z0-9]+-[0-9]{2}", public_error_code)
            else "unknown"
        )
        super().__init__(f"{safe_case} failed with public code {safe_code}")
        self.case_name = safe_case
        self.public_error_code = safe_code


def require(response: httpx.Response, expected: int = 200) -> dict[str, Any]:
    if response.status_code != expected:
        raise RuntimeError(
            f"Unexpected local API status {response.status_code} at {response.request.url.path}"
        )
    payload = response.json()
    if not isinstance(payload, dict):
        raise RuntimeError("Local API did not return an object")
    return payload


def _run_media_command(command: list[str]) -> None:
    try:
        subprocess.run(
            command,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        raise RuntimeError("Synthetic media generation failed") from exc


def _font_path() -> Path:
    candidates = (
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf"),
    )
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    raise RuntimeError("A Unicode font is required for the format rehearsal")


def _create_native_documents(directory: Path) -> list[tuple[Path, str]]:
    lesson = (
        "Fotosentez, bitkilerin isik enerjisini kimyasal enerjiye donusturdugu "
        "temel biyolojik surectir. Kloroplasttaki isiga bagli tepkimeler ATP ve "
        "NADPH uretir. Calvin dongusu karbondioksiti kullanarak karbonhidrat "
        "sentezler. Stomalar gaz alisverisini duzenler; su, sicaklik ve isik "
        "siddeti surecin hizini etkiler. "
    )

    text_path = directory / "native-notes.txt"
    text_path.write_text((lesson + "\n") * 12, encoding="utf-8")

    markdown_path = directory / "native-outline.md"
    markdown_path.write_text(
        "# Fotosentez\n\n## Isiga bagli tepkimeler\n\nATP ve NADPH uretilir.\n",
        encoding="utf-8",
    )

    docx_path = directory / "native-workbook.docx"
    document = Document()
    document.add_heading("Fotosentez calisma notu", level=1)
    document.add_paragraph("Kloroplastlar isik enerjisini kimyasal enerjiye donusturur.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Asama"
    table.cell(0, 1).text = "Urun"
    document.save(docx_path)

    pdf_path = directory / "native-handout.pdf"
    font_path = _font_path()
    pdfmetrics.registerFont(TTFont("RehearsalSans", str(font_path)))
    pdf = canvas.Canvas(str(pdf_path))
    pdf.setFont("RehearsalSans", 15)
    pdf.drawString(54, 780, "Fotosentez ve enerji donusumu")
    pdf.setFont("RehearsalSans", 11)
    pdf.drawString(54, 750, "Klorofil isigi soğurur; ATP ve NADPH kimyasal enerji tasir.")
    pdf.drawString(54, 730, "Calvin dongusu karbondioksitten organik molekuller uretir.")
    pdf.save()

    pptx_path = directory / "native-slides.pptx"
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Fotosentezin asamalari"
    first.placeholders[1].text = "Isiga bagli tepkimeler ATP ve NADPH uretir."
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Calvin dongusu"
    second.placeholders[1].text = "Karbondioksit organik molekullere donusturulur."
    slide_table = second.shapes.add_table(1, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
    slide_table.cell(0, 0).text = "PPTX TABLE SENTINEL"
    slide_table.cell(0, 1).text = "Karbondioksit"
    second.notes_slide.notes_text_frame.text = (
        "PPTX NOTES SENTINEL: Calvin dongusunun kloroplast stromasinda "
        "oldugunu ayrintili olarak acikla."
    )

    # Exercise the real picture-placeholder path and prove that long speaker
    # notes do not suppress OCR of an otherwise image-only slide.
    pptx_scan_path = directory / "pptx-placeholder-scan.png"
    pptx_scan = Image.new("RGB", (1800, 1000), "white")
    pptx_draw = ImageDraw.Draw(pptx_scan)
    pptx_title_font = ImageFont.truetype(str(font_path), 74)
    pptx_body_font = ImageFont.truetype(str(font_path), 52)
    pptx_draw.text((140, 210), "PPTX OCR REHEARSAL", fill="black", font=pptx_title_font)
    pptx_draw.text(
        (140, 390),
        "CHLOROPLAST ENERGY TEST",
        fill="black",
        font=pptx_body_font,
    )
    pptx_scan.save(pptx_scan_path, format="PNG", dpi=(200, 200))
    scanned = presentation.slides.add_slide(presentation.slide_layouts[8])
    scanned.placeholders[1].insert_picture(str(pptx_scan_path))
    scanned.notes_slide.notes_text_frame.text = (
        "This intentionally long speaker note contains more than the native text "
        "threshold and must never prevent OCR of the visual slide content."
    )
    presentation.save(pptx_path)

    return [
        (text_path, "text/plain"),
        (markdown_path, "text/markdown"),
        (
            docx_path,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
        (pdf_path, "application/pdf"),
        (
            pptx_path,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
    ]


def _create_ocr_images(directory: Path) -> list[tuple[Path, str]]:
    image_path = directory / "ocr-scan-source.png"
    image = Image.new("RGB", (1800, 1050), "white")
    draw = ImageDraw.Draw(image)
    font_path = _font_path()
    title_font = ImageFont.truetype(str(font_path), 64)
    body_font = ImageFont.truetype(str(font_path), 46)
    draw.text((110, 95), "FOTOSENTEZ DERS NOTU", fill="black", font=title_font)
    rows = (
        "Bitkiler isik enerjisini kimyasal enerjiye donusturur.",
        "Kloroplastlarda ATP ve NADPH molekulleri uretilir.",
        "Calvin dongusu karbondioksiti kullanarak seker sentezler.",
        "Su miktari ve isik siddeti fotosentez hizini etkiler.",
    )
    for index, row in enumerate(rows):
        draw.text((110, 255 + index * 155), row, fill="black", font=body_font)
    image.save(image_path, format="PNG", optimize=True, dpi=(200, 200))
    variants = (
        ("png", "image/png", "PNG"),
        ("jpg", "image/jpeg", "JPEG"),
        ("jpeg", "image/jpeg", "JPEG"),
        ("webp", "image/webp", "WEBP"),
        ("tif", "image/tiff", "TIFF"),
        ("tiff", "image/tiff", "TIFF"),
    )
    outputs: list[tuple[Path, str]] = []
    for suffix, media_type, image_format in variants:
        output = directory / f"ocr-scan.{suffix}"
        image.save(output, format=image_format, dpi=(200, 200))
        outputs.append((output, media_type))
    return outputs


def _prepared_synthetic_speech(directory: Path) -> tuple[Path, str]:
    prepared = Path("/tmp/lecturesift-rehearsal-synthetic-lecture.mp3")
    if os.getenv("LECTURESIFT_REHEARSAL_AI_PROVIDER") != "dedicated":
        raise RuntimeError("dedicated rehearsal AI source was not enabled")
    if not prepared.is_file() or prepared.is_symlink() or prepared.stat().st_size < 1024:
        raise RuntimeError("dedicated rehearsal synthetic speech source is missing")
    audio_path = directory / "synthetic-lecture.mp3"
    shutil.copyfile(prepared, audio_path)
    if not audio_path.is_file() or audio_path.stat().st_size < 1024:
        raise RuntimeError("Synthetic MP3 was not created")
    return audio_path, "audio/mpeg"


def _create_eight_second_video(directory: Path, audio_path: Path) -> tuple[Path, str]:
    slide_path = directory / "video-slide.png"
    image = Image.new("RGB", (960, 540), (15, 35, 66))
    draw = ImageDraw.Draw(image)
    font_path = _font_path()
    title_font = ImageFont.truetype(str(font_path), 52)
    body_font = ImageFont.truetype(str(font_path), 34)
    draw.text((75, 120), "LectureSift Rehearsal", fill=(255, 255, 255), font=title_font)
    draw.text((75, 225), "Fotosentez ve enerji donusumu", fill=(118, 210, 255), font=body_font)
    draw.text((75, 295), "8 saniyelik sentetik ders videosu", fill=(220, 230, 242), font=body_font)
    image.save(slide_path, format="PNG", optimize=True)

    video_path = directory / "synthetic-lecture-8s.mp4"
    _run_media_command(
        [
            "ffmpeg",
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            "-loop",
            "1",
            "-i",
            str(slide_path),
            "-i",
            str(audio_path),
            "-filter_complex",
            "[1:a]apad=pad_dur=8[audio]",
            "-map",
            "0:v:0",
            "-map",
            "[audio]",
            "-t",
            "8",
            "-r",
            "12",
            "-c:v",
            "libx264",
            "-preset",
            "veryfast",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            "-movflags",
            "+faststart",
            str(video_path),
        ]
    )
    try:
        probe = subprocess.run(
            [
                "ffprobe",
                "-v",
                "error",
                "-show_entries",
                "format=duration",
                "-of",
                "default=noprint_wrappers=1:nokey=1",
                str(video_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=30,
        )
        duration = float(probe.stdout.strip())
    except (OSError, ValueError, subprocess.SubprocessError) as exc:
        raise RuntimeError("Synthetic MP4 duration could not be verified") from exc
    if not 7.9 <= duration <= 8.1:
        raise RuntimeError("Synthetic MP4 is not eight seconds long")
    return video_path, "video/mp4"


def _job_data(*, include_slides: bool) -> dict[str, str]:
    return {
        "source_language": "tr",
        "output_language": "tr",
        "summary_style": "standard",
        "quiz_count": "0",
        "flashcard_count": "0",
        "translate_transcript": "false",
        "output_formats": "pdf",
        "job_type": "study_pack",
        # The lightweight rehearsal separately proves real study-pack model
        # output. This matrix isolates ingestion/OCR/transcription/export so
        # one run cannot spend four redundant summary completions.
        "include_summary": "false",
        "include_transcript": "true",
        "include_slides": "true" if include_slides else "false",
        "transcript_timestamps": "false",
        "speaker_detection": "false",
    }


def _submit_job(
    client: httpx.Client,
    auth: dict[str, str],
    sources: list[tuple[Path, str]],
    *,
    case_name: str,
    include_slides: bool,
) -> str:
    multipart = [
        ("files", (path.name, path.read_bytes(), media_type))
        for path, media_type in sources
    ]
    response = client.post(
        "/jobs",
        headers=auth,
        files=multipart,
        data=_job_data(include_slides=include_slides),
        timeout=120.0,
    )
    if response.status_code != 200:
        try:
            detail = response.json().get("detail") or {}
            public_code = str(detail.get("code") or "unknown")
        except (AttributeError, TypeError, ValueError):
            public_code = "unknown"
        raise FormatRehearsalFailure(case_name, public_code)
    created = require(response)
    job_id = str(created.get("job_id") or "")
    if not job_id:
        raise RuntimeError("Local API did not create a job identifier")
    return job_id


def _wait_for_durable_jobs(
    client: httpx.Client,
    auth: dict[str, str],
    jobs: dict[str, str],
) -> dict[str, dict[str, Any]]:
    terminal: dict[str, dict[str, Any]] = {}
    deadline = time.monotonic() + TERMINAL_TIMEOUT_SECONDS
    while len(terminal) < len(jobs) and time.monotonic() < deadline:
        for case_name, job_id in jobs.items():
            if case_name in terminal:
                continue
            job = require(client.get(f"/jobs/{job_id}", headers=auth))
            status = str(job.get("status") or "")
            if status == "error":
                error_code = str(job.get("error_code") or "unknown")
                raise FormatRehearsalFailure(case_name, error_code)
            if status != "done":
                continue
            if not (
                str(job.get("worker_state") or "") == "done"
                and str(job.get("queue_mode") or "") == "celery"
                and int(job.get("percent") or 0) == 100
                and str(job.get("stage") or "") == "done"
            ):
                raise RuntimeError(f"{case_name} bypassed a durable Celery completion gate")
            terminal[case_name] = job
        if len(terminal) < len(jobs):
            time.sleep(2)
    if len(terminal) != len(jobs):
        raise RuntimeError("Timed out while waiting for durable format jobs")
    return terminal


def _has_pdf_artifact(result: dict[str, Any]) -> bool:
    return any(
        str(item.get("format") or "").casefold() == "pdf"
        for item in result.get("artifacts", [])
        if isinstance(item, dict)
    )


def _assert_result(case_name: str, job_id: str, result: dict[str, Any]) -> None:
    if result.get("job_id") != job_id:
        raise RuntimeError(f"{case_name} returned another job's result")
    if not str(result.get("transcript") or "").strip():
        raise RuntimeError(f"{case_name} did not produce source text")
    if not _has_pdf_artifact(result):
        raise RuntimeError(f"{case_name} did not publish its PDF artifact manifest")

    sources = result.get("sources") or {}
    diagnostics = result.get("diagnostics") or {}
    if case_name == "native_documents":
        documents = [item for item in sources.get("documents", []) if isinstance(item, dict)]
        document_types = {str(item.get("type") or "").casefold() for item in documents}
        pptx_document = next(
            (item for item in documents if str(item.get("type") or "").casefold() == "pptx"),
            {},
        )
        transcript = str(result.get("transcript") or "").casefold()
        if (
            sources.get("mode") != "documents"
            or len(documents) != 5
            or document_types != {"txt", "md", "docx", "pdf", "pptx"}
            or any(int(item.get("characters") or 0) <= 0 for item in documents)
            or diagnostics.get("ocr_used") is not True
            or int(diagnostics.get("ocr_pages") or 0) < 1
            or pptx_document.get("ocr_used") is not True
            or int(pptx_document.get("ocr_pages") or 0) != 1
            or int(pptx_document.get("ocr_images") or 0) != 1
            or "pptx table sentinel" not in transcript
            or "pptx notes sentinel" not in transcript
            or "pptx ocr rehearsal" not in transcript
        ):
            raise RuntimeError("Native document batch coverage was incomplete")
    elif case_name == "ocr_images":
        documents = [item for item in sources.get("documents", []) if isinstance(item, dict)]
        document_types = {str(item.get("type") or "").casefold() for item in documents}
        if (
            len(documents) != 6
            or document_types != {"png", "jpg", "jpeg", "webp", "tif", "tiff"}
            or any(int(item.get("characters") or 0) <= 0 for item in documents)
            or diagnostics.get("ocr_used") is not True
            or int(diagnostics.get("ocr_pages") or 0) < 6
            or "fotosentez" not in str(result.get("transcript") or "").casefold()
        ):
            raise RuntimeError("Image OCR coverage was incomplete")
    elif case_name == "mp3_audio":
        audio_files = [str(value).casefold() for value in sources.get("audio_files", [])]
        if (
            not audio_files
            or not all(value.endswith(".mp3") for value in audio_files)
            or "fotosentez" not in str(result.get("transcript") or "").casefold()
        ):
            raise RuntimeError("MP3 source coverage was incomplete")
    elif case_name == "mp4_video":
        audio_files = [str(value).casefold() for value in sources.get("audio_files", [])]
        visual_files = [str(value).casefold() for value in sources.get("visual_files", [])]
        if (
            not audio_files
            or not visual_files
            or not all(value.endswith(".mp4") for value in audio_files + visual_files)
            or int(diagnostics.get("final_unique_slides") or 0) < 1
            or not result.get("slides")
            or "fotosentez" not in str(result.get("transcript") or "").casefold()
        ):
            raise RuntimeError("MP4 audio/visual coverage was incomplete")


def _verify_mp3(path: Path) -> None:
    if path.is_symlink() or not path.is_file() or path.stat().st_size <= 0:
        raise RuntimeError("R2 MP3 artifact was empty or unsafe")
    completed = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-select_streams",
            "a:0",
            "-show_entries",
            "stream=codec_name",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(path),
        ],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=30,
    )
    if completed.returncode != 0 or completed.stdout.strip().casefold() != "mp3":
        raise RuntimeError("R2 MP3 artifact did not contain a valid MP3 audio stream")


_WINDOWS_RESERVED_NAMES = {
    "con",
    "prn",
    "aux",
    "nul",
    *(f"com{number}" for number in range(1, 10)),
    *(f"lpt{number}" for number in range(1, 10)),
}
_MAX_REHEARSAL_ZIP_MEMBERS = 64
_MAX_REHEARSAL_ZIP_MEMBER_BYTES = 64 * 1024 * 1024
_MAX_REHEARSAL_ZIP_TOTAL_BYTES = 128 * 1024 * 1024


def _safe_leaf_name(value: object, *, label: str) -> str:
    """Reject names that could escape or alias on POSIX or Windows clients."""

    if not isinstance(value, str) or not value:
        raise RuntimeError(f"{label} was missing or was not a string")
    name = value
    posix_name = PurePosixPath(name)
    windows_name = PureWindowsPath(name)
    windows_stem = name.split(".", 1)[0].rstrip(" .").casefold()
    if (
        name in {".", ".."}
        or "/" in name
        or "\\" in name
        or posix_name.is_absolute()
        or windows_name.is_absolute()
        or posix_name.name != name
        or windows_name.name != name
        or bool(windows_name.drive)
        or bool(windows_name.root)
        or name.endswith((" ", "."))
        or windows_stem in _WINDOWS_RESERVED_NAMES
        or any(character in '<>:"|?*' for character in name)
        or any(ord(character) < 32 or ord(character) == 127 for character in name)
    ):
        raise RuntimeError(f"{label} was not a safe cross-platform leaf name")
    return name


def _has_casefold_duplicates(names: list[str]) -> bool:
    return len(names) != len({name.casefold() for name in names})


def _verify_r2_payloads(
    case_name: str,
    job_id: str,
    result: dict[str, Any],
    local_dir: Path,
    archive_key: str,
) -> list[str]:
    raw_artifacts = result.get("artifacts")
    if not isinstance(raw_artifacts, list) or any(
        not isinstance(item, dict) for item in raw_artifacts
    ):
        raise RuntimeError("R2 artifact manifest was missing or malformed")
    artifacts = raw_artifacts
    artifact_names = [
        _safe_leaf_name(item.get("file"), label="R2 artifact filename")
        for item in artifacts
    ]
    raw_artifact_sizes = [item.get("size_bytes") for item in artifacts]
    if any(
        isinstance(size, bool) or not isinstance(size, int)
        for size in raw_artifact_sizes
    ):
        raise RuntimeError("R2 artifact manifest contained an invalid size")
    artifact_sizes = raw_artifact_sizes
    if (
        not artifact_names
        or len(artifact_names) > _MAX_REHEARSAL_ZIP_MEMBERS
        or len(artifact_names) != len(set(artifact_names))
        or _has_casefold_duplicates(artifact_names)
        or any(
            size <= 0 or size > _MAX_REHEARSAL_ZIP_MEMBER_BYTES
            for size in artifact_sizes
        )
        or sum(artifact_sizes) > _MAX_REHEARSAL_ZIP_TOTAL_BYTES
    ):
        raise RuntimeError("R2 artifact manifest contained an unsafe or empty file")
    pdf_artifact = next(
        (
            item
            for item in artifacts
            if str(item.get("format") or "").casefold() == "pdf"
        ),
        None,
    )
    if not pdf_artifact:
        raise RuntimeError("R2 payload probe could not identify a PDF artifact")
    pdf_name = str(pdf_artifact["file"])
    pdf_key = f"jobs/{job_id}/package/{pdf_name}"

    prefix = f"jobs/{job_id}/"
    if not archive_key.startswith(prefix):
        raise RuntimeError("R2 result did not bind its ZIP archive to the owning job")
    archive_relative = archive_key[len(prefix) :]
    archive_relative_path = PurePosixPath(archive_relative)
    archive_parts = archive_relative.split("/")
    if (
        not archive_relative
        or "\\" in archive_relative
        or archive_relative_path.is_absolute()
        or ".." in archive_relative_path.parts
        or any(part in {"", "."} for part in archive_parts)
        or any(ord(character) < 32 or ord(character) == 127 for character in archive_relative)
        or archive_relative_path.suffix.casefold() != ".zip"
    ):
        raise RuntimeError("R2 result exposed an unsafe ZIP archive path")
    for archive_part in archive_parts:
        _safe_leaf_name(archive_part, label="R2 ZIP archive path segment")

    media_case = case_name in {"mp3_audio", "mp4_video"}
    audio_artifacts = [
        item
        for item in artifacts
        if str(item.get("format") or "").casefold() == "mp3"
    ]
    if (media_case and len(audio_artifacts) != 1) or (
        not media_case and audio_artifacts
    ):
        raise RuntimeError("R2 media artifact manifest did not match its format case")

    evidence = ["pdf_sample", "archive_zip"]
    keys = [pdf_key, archive_key]
    audio_name = ""
    if media_case:
        audio_name = str(audio_artifacts[0]["file"])
        keys.append(f"jobs/{job_id}/package/{audio_name}")
        evidence.append("audio_mp3")

    raw_slides = result.get("slides", [])
    if not isinstance(raw_slides, list) or any(
        not isinstance(item, dict) for item in raw_slides
    ):
        raise RuntimeError("R2 slide manifest was malformed")
    slides = raw_slides
    if case_name == "mp4_video" and not slides:
        raise RuntimeError("R2 video result did not expose a slide sample")
    if case_name != "mp4_video" and slides:
        raise RuntimeError("R2 non-video result unexpectedly exposed slide payloads")
    slide_name = ""
    if case_name == "mp4_video":
        slide_names = [
            _safe_leaf_name(item.get("file"), label="R2 slide filename")
            for item in slides
        ]
        if _has_casefold_duplicates(slide_names):
            raise RuntimeError("R2 slide manifest contained aliased filenames")
        slide_name = slide_names[0]
        keys.append(f"jobs/{job_id}/slides/{slide_name}")
        evidence.append("slide_sample")

    shutil.rmtree(local_dir / "package", ignore_errors=True)
    archive_local = local_dir.joinpath(*archive_relative_path.parts)
    archive_local.unlink(missing_ok=True)
    if STORAGE.materialize_files(job_id, local_dir, keys) != len(keys):
        raise RuntimeError("Expected R2 payloads could not all be materialized")

    pdf_path = local_dir / "package" / pdf_name
    if not pdf_path.read_bytes().startswith(b"%PDF-"):
        raise RuntimeError("R2 PDF artifact signature was invalid")

    audio_path = local_dir / "package" / audio_name if audio_name else None
    if audio_path is not None:
        _verify_mp3(audio_path)

    if slide_name:
        with Image.open(local_dir / "slides" / slide_name) as slide_image:
            slide_image.verify()
            if slide_image.format != "JPEG":
                raise RuntimeError("R2 slide image was not a valid JPEG")

    try:
        with zipfile.ZipFile(archive_local) as archive:
            members = archive.infolist()
            if len(members) > _MAX_REHEARSAL_ZIP_MEMBERS:
                raise RuntimeError("R2 ZIP archive exceeded the rehearsal member limit")
            names = [
                _safe_leaf_name(member.filename, label="R2 ZIP member filename")
                for member in members
            ]
            total_uncompressed_bytes = sum(member.file_size for member in members)
            if (
                not names
                or len(names) != len(set(names))
                or _has_casefold_duplicates(names)
                or set(names) != set(artifact_names)
                or total_uncompressed_bytes > _MAX_REHEARSAL_ZIP_TOTAL_BYTES
                or any(
                    member.is_dir()
                    or member.file_size < 0
                    or member.file_size > _MAX_REHEARSAL_ZIP_MEMBER_BYTES
                    or stat.S_ISLNK(member.external_attr >> 16)
                    or stat.S_IFMT(member.external_attr >> 16)
                    not in {0, stat.S_IFREG}
                    for member in members
                )
            ):
                raise RuntimeError("R2 ZIP archive manifest was unsafe or incomplete")
            if archive.testzip() is not None:
                raise RuntimeError("R2 ZIP archive failed its CRC check")
            for artifact in artifacts:
                name = str(artifact["file"])
                payload = archive.read(name)
                if len(payload) != int(artifact["size_bytes"]):
                    raise RuntimeError("R2 ZIP member size did not match its manifest")
                if str(artifact.get("format") or "").casefold() == "pdf" and not payload.startswith(
                    b"%PDF-"
                ):
                    raise RuntimeError("R2 ZIP contained an invalid PDF member")
            if audio_path is not None and archive.read(audio_name) != audio_path.read_bytes():
                raise RuntimeError("R2 ZIP did not contain the verified MP3 payload")
    except (OSError, zipfile.BadZipFile) as exc:
        raise RuntimeError("R2 ZIP archive could not be verified") from exc

    return evidence


def _job_is_durably_terminal(job: dict[str, Any]) -> bool:
    """Return true only when cleanup cannot race a Celery publication."""
    status = str(job.get("status") or "")
    queue_mode = str(job.get("queue_mode") or "")
    worker_state = str(job.get("worker_state") or "")
    if queue_mode == "celery":
        if status == "done":
            return worker_state == "done"
        if status == "error":
            return worker_state in {"failed", "rejected", "unavailable"}
        return False
    return status in {"done", "error"}


def main() -> None:
    global CURRENT_STAGE

    CURRENT_STAGE = "environment_guard"
    database_name = make_url(config.DATABASE_URL).database or ""
    if os.getenv("LECTURESIFT_REHEARSAL") != "1" or not database_name.startswith(
        "lecturesift_rehearsal_"
    ):
        raise RuntimeError("Refusing to run outside an explicit rehearsal database")
    if not STORAGE.remote:
        raise RuntimeError("Private object storage is required for this rehearsal")

    run_id = uuid.uuid4().hex
    test_email = f"format-rehearsal-{run_id}@example.invalid"
    password = secrets.token_urlsafe(36)
    user_id = ""
    account_closed = False
    cleanup_complete = False
    jobs: dict[str, str] = {}
    job_dirs: dict[str, Path] = {}
    available_cases = {"native_documents", "ocr_images", "mp3_audio", "mp4_video"}
    requested_cases = {
        value.strip()
        for value in os.getenv("LECTURESIFT_REHEARSAL_FORMAT_CASES", "").split(",")
        if value.strip()
    }
    selected_cases = requested_cases or available_cases
    if not selected_cases or not selected_cases.issubset(available_cases):
        raise RuntimeError("Unknown rehearsal format case selection")
    ai_provider = os.getenv("LECTURESIFT_REHEARSAL_AI_PROVIDER", "")
    if ai_provider not in {"dedicated", "intentionally_absent"}:
        raise RuntimeError("Rehearsal AI provider state is not explicit")
    ai_cases = {"mp3_audio", "mp4_video"}.intersection(selected_cases)
    skipped_cases: dict[str, str] = {}
    if ai_provider == "intentionally_absent":
        skipped_cases = {
            case: "dedicated_rehearsal_openai_key_absent" for case in sorted(ai_cases)
        }
        selected_cases = selected_cases.difference(ai_cases)
    if not selected_cases:
        raise RuntimeError("All requested cases require an absent dedicated AI provider")

    def best_effort_cleanup() -> None:
        nonlocal account_closed, cleanup_complete
        if cleanup_complete:
            return
        # Never remove metadata or R2 sources while a worker can still be
        # processing them. A timed-out rehearsal is left isolated for manual
        # inspection instead of creating an orphan-publication race.
        try:
            active_jobs = [JOBS.get(job_id) or {} for job_id in jobs.values()]
        except Exception:
            return
        if any(not _job_is_durably_terminal(job) for job in active_jobs):
            return
        if user_id and not account_closed:
            try:
                admin_close_user_account(
                    user_id,
                    confirmation_email=test_email,
                    reason="Migration format rehearsal cleanup",
                    actor="format-rehearsal",
                )
                account_closed = True
            except Exception:
                pass
        if user_id:
            try:
                JOBS.delete_for_user(user_id)
            except Exception:
                pass
        for job_id in jobs.values():
            try:
                STORAGE.delete_job(job_id)
            except Exception:
                pass
            shutil.rmtree(config.WORK_DIR / job_id, ignore_errors=True)

    atexit.register(best_effort_cleanup)

    try:
        CURRENT_STAGE = "account"
        registration = register_user(
            test_email,
            password,
            "Format",
            "Rehearsal",
            country_code="TR",
        )
        account = verify_email(registration["verification_token"])
        user_id = str(account["user"]["id"])
        admin_set_user_subscription(
            user_id,
            plan_code="pro",
            interval="monthly",
            duration_days=1,
            actor="format-rehearsal",
        )
        token = str(account["token"])
        auth = {"Authorization": f"Bearer {token}"}
        admin_auth = {"Authorization": f"Bearer {config.ADMIN_ADMIN}"}

        with tempfile.TemporaryDirectory(prefix="lecturesift-formats-rehearsal-") as raw_dir:
            source_dir = Path(raw_dir)
            CURRENT_STAGE = "source_generation"
            cases: dict[str, tuple[list[tuple[Path, str]], bool]] = {}
            if "native_documents" in selected_cases:
                cases["native_documents"] = (_create_native_documents(source_dir), False)
            if "ocr_images" in selected_cases:
                cases["ocr_images"] = (_create_ocr_images(source_dir), False)
            if selected_cases.intersection({"mp3_audio", "mp4_video"}):
                speech_audio = _prepared_synthetic_speech(source_dir)
                if "mp3_audio" in selected_cases:
                    cases["mp3_audio"] = ([speech_audio], False)
                if "mp4_video" in selected_cases:
                    video = _create_eight_second_video(source_dir, speech_audio[0])
                    cases["mp4_video"] = ([video], True)

            with httpx.Client(base_url=BASE_URL, timeout=30.0) as client:
                CURRENT_STAGE = "job_submission"
                for case_name, (sources, include_slides) in cases.items():
                    jobs[case_name] = _submit_job(
                        client,
                        auth,
                        sources,
                        case_name=case_name,
                        include_slides=include_slides,
                    )
                    job_dirs[case_name] = config.WORK_DIR / jobs[case_name]

                CURRENT_STAGE = "celery_completion"
                started = time.monotonic()
                _wait_for_durable_jobs(client, auth, jobs)
                elapsed_seconds = round(time.monotonic() - started, 2)

                CURRENT_STAGE = "r2_reopen"
                reopened = 0
                payloads_verified = 0
                payload_evidence_by_case: dict[str, list[str]] = {}
                for case_name, job_id in jobs.items():
                    persisted = JOBS.get(job_id) or {}
                    if not (
                        str(persisted.get("remote_prefix") or "") == f"jobs/{job_id}/"
                        and str(persisted.get("remote_result_key") or "")
                        == f"jobs/{job_id}/result.json"
                        and int(persisted.get("remote_file_count") or 0) >= 2
                    ):
                        raise RuntimeError(f"{case_name} did not publish a durable R2 result")

                    local_dir = job_dirs[case_name]
                    shutil.rmtree(local_dir, ignore_errors=True)
                    if local_dir.exists():
                        raise RuntimeError("Local eviction failed before the R2 reopen probe")
                    result = require(client.get(f"/jobs/{job_id}/result", headers=auth))
                    _assert_result(case_name, job_id, result)
                    evidence = _verify_r2_payloads(
                        case_name,
                        job_id,
                        result,
                        local_dir,
                        str(persisted.get("remote_download_key") or ""),
                    )
                    payload_evidence_by_case[case_name] = evidence
                    payloads_verified += len(evidence)
                    if not (local_dir / "result.json").is_file():
                        raise RuntimeError(f"{case_name} result was not reopened from R2")
                    reopened += 1

                CURRENT_STAGE = "account_cleanup"
                closed = require(
                    client.request(
                        "DELETE",
                        f"/billing/admin/users/{user_id}",
                        headers=admin_auth,
                        json={
                            "confirmation_email": test_email,
                            "reason": "Migration format rehearsal cleanup",
                        },
                    )
                )
                account_closed = bool(closed.get("ok"))
                if not account_closed or int(closed.get("deleted_jobs") or 0) != len(jobs):
                    raise RuntimeError("Rehearsal account cleanup did not remove every job")

        CURRENT_STAGE = "cleanup_verification"
        if any(JOBS.get(job_id) is not None for job_id in jobs.values()):
            raise RuntimeError("Rehearsal job metadata remained after account closure")
        if any(path.exists() for path in job_dirs.values()):
            raise RuntimeError("Rehearsal local job data remained after account closure")
        residual_objects = sum(STORAGE.delete_job(job_id) for job_id in jobs.values())
        if residual_objects:
            raise RuntimeError("Rehearsal R2 objects remained after account closure")

        cleanup_complete = True
        CURRENT_STAGE = "complete"
        print(
            json.dumps(
                {
                    "ok": True,
                    "cases": sorted(cases),
                    "requested_cases": sorted(requested_cases or available_cases),
                    "skipped_cases": skipped_cases,
                    "ai_provider_tested": ai_provider == "dedicated",
                    "ai_provider_state": ai_provider,
                    "formats": sorted(
                        format_name
                        for case_name in cases
                        for format_name in {
                            "native_documents": {"txt", "md", "docx", "pdf", "pptx"},
                            "ocr_images": {"png", "jpg", "jpeg", "webp", "tif", "tiff"},
                            "mp3_audio": {"mp3"},
                            "mp4_video": {"mp4_8s"},
                        }[case_name]
                    ),
                    "jobs_submitted": len(jobs),
                    "celery_jobs_done": len(jobs),
                    "r2_results_reopened": reopened,
                    "r2_payloads_verified": payloads_verified,
                    "r2_payloads_verified_by_case": payload_evidence_by_case,
                    "job_metadata_removed": len(jobs),
                    "r2_residual_objects": residual_objects,
                    "rehearsal_account_closed": account_closed,
                    "rehearsal_user_id": user_id,
                    "rehearsal_email": test_email,
                    "rehearsal_job_ids": sorted(jobs.values()),
                    "elapsed_seconds": elapsed_seconds,
                },
                ensure_ascii=False,
                sort_keys=True,
            )
        )
    finally:
        best_effort_cleanup()


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        safe_failure = (
            {
                "failed_case": exc.case_name,
                "public_error_code": exc.public_error_code,
            }
            if isinstance(exc, FormatRehearsalFailure)
            else {}
        )
        print(
            json.dumps(
                {
                    "ok": False,
                    "failed_stage": CURRENT_STAGE,
                    "error_type": type(exc).__name__,
                    **safe_failure,
                },
                sort_keys=True,
            )
        )
        raise SystemExit(1) from None
