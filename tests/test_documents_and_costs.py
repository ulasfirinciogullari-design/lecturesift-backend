from __future__ import annotations

import json
import os
import shutil
import threading
import time
import uuid
from datetime import date, datetime, timezone
from pathlib import Path
from types import SimpleNamespace

import pytest
from docx import Document
from fastapi.testclient import TestClient
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from reportlab.pdfgen import canvas

import lecturesift.app as app_module
import lecturesift.costs as costs
import lecturesift.documents as document_service
from lecturesift import config
from lecturesift.app import app
from lecturesift.billing_service import (
    IYZICO_BANK_TRANSFER_PROVIDER,
    PAYMENT_ORDERS,
    create_payment_order,
    register_user,
    verify_email,
)
from lecturesift.documents import extract_documents
from lecturesift.durable_runtime import _preflight_documents
from lecturesift.errors import LectureSiftError
from lecturesift.jobs import JOBS
from lecturesift.rollout_routes import install_rollout_routes


def _account_headers() -> dict[str, str]:
    created = register_user(
        f"document-{uuid.uuid4()}@example.com",
        "Strong-test-password1",
        "Document",
        "User",
    )
    token = verify_email(created["verification_token"])["token"]
    return {"Authorization": f"Bearer {token}"}


def test_text_docx_pptx_and_pdf_extraction(tmp_path: Path):
    text_path = tmp_path / "notes.txt"
    text_path.write_text(" ".join(f"word{index}" for index in range(205)), encoding="utf-8")

    docx_path = tmp_path / "lesson.docx"
    document = Document()
    document.add_heading("Energy", level=1)
    document.add_paragraph("Potential and kinetic energy are related.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Unit"
    table.cell(0, 1).text = "Joule"
    document.save(docx_path)

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Conservation law"
    slide.placeholders[1].text = "Energy changes form but is conserved."
    slide_table = slide.shapes.add_table(1, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
    slide_table.cell(0, 0).text = "Equation"
    slide_table.cell(0, 1).text = "E = K + U"
    slide.notes_slide.notes_text_frame.text = "Mention isolated systems and friction."
    presentation.save(pptx_path)

    pdf_path = tmp_path / "handout.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 760, "Work transfers energy between systems.")
    pdf.save()

    result = extract_documents([text_path, docx_path, pptx_path, pdf_path])
    assert result["credit_minutes"] >= 2
    assert result["words"] >= 205
    assert [item["type"] for item in result["documents"]] == ["txt", "docx", "pptx", "pdf"]
    assert "Potential and kinetic energy" in result["text"]
    assert "Conservation law" in result["text"]
    assert "Equation | E = K + U" in result["text"]
    assert "Mention isolated systems and friction" in result["text"]
    assert "Work transfers energy" in result["text"]


def test_image_only_pptx_is_read_with_ocr_and_preflight_can_defer(tmp_path: Path, monkeypatch):
    image_path = tmp_path / "slide-scan.png"
    Image.new("RGB", (1600, 900), "white").save(image_path)

    pptx_path = tmp_path / "scanned-slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width)
    presentation.save(pptx_path)

    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda image, source_language="auto": (
            "Fotosentez ışık enerjisini kimyasal enerjiye dönüştürür.",
            "tur+eng",
        ),
    )

    preflight = extract_documents(
        [pptx_path],
        source_language="tr",
        enable_ocr=False,
        allow_ocr_pending=True,
    )
    assert preflight["documents"][0]["ocr_required"] is True
    assert preflight["documents"][0]["ocr_pages"] == 1

    result = extract_documents([pptx_path], source_language="tr")
    assert "Fotosentez" in result["text"]
    assert result["ocr_used"] is True
    assert result["ocr_pages"] == 1
    assert result["documents"][0]["ocr_language"] == "tur+eng"


def test_picture_placeholder_pptx_reaches_ocr(tmp_path: Path, monkeypatch):
    image_path = tmp_path / "placeholder-scan.png"
    Image.new("RGB", (1600, 900), "white").save(image_path)

    pptx_path = tmp_path / "picture-placeholder.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    inserted = slide.placeholders[1].insert_picture(str(image_path))
    assert inserted.shape_type != document_service.MSO_SHAPE_TYPE.PICTURE
    presentation.save(pptx_path)

    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda image, source_language="auto": ("Placeholder OCR content", "eng"),
    )

    result = extract_documents([pptx_path], source_language="en")

    assert "Placeholder OCR content" in result["text"]
    assert result["documents"][0]["ocr_pages"] == 1
    assert result["documents"][0]["ocr_images"] == 1
    assert result["documents"][0]["native_text_pages"] == 0


def test_pptx_ocrs_all_substantial_images_and_notes_do_not_suppress_it(tmp_path: Path, monkeypatch):
    left_path = tmp_path / "left.png"
    right_path = tmp_path / "right.png"
    Image.new("RGB", (900, 900), "red").save(left_path)
    Image.new("RGB", (900, 900), "blue").save(right_path)

    pptx_path = tmp_path / "two-scans-with-notes.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(left_path), Inches(0.5), Inches(1), width=Inches(5.5))
    slide.shapes.add_picture(str(right_path), Inches(7.2), Inches(1), width=Inches(5.5))
    slide.notes_slide.notes_text_frame.text = (
        "This deliberately long speaker note exceeds the native-text threshold "
        "but is not visible slide content and therefore must not suppress OCR."
    )
    presentation.save(pptx_path)

    state = {"active": 0, "maximum": 0}
    lock = threading.Lock()
    gate = threading.Barrier(2, timeout=3)

    def fake_ocr(image, source_language="auto"):
        with lock:
            state["active"] += 1
            state["maximum"] = max(state["maximum"], state["active"])
        gate.wait()
        pixel = image.convert("RGB").getpixel((0, 0))
        text = "LEFT OCR PANEL" if pixel[0] > pixel[2] else "RIGHT OCR PANEL"
        with lock:
            state["active"] -= 1
        return text, "eng"

    monkeypatch.setattr(document_service, "OCR_PARALLELISM", 2)
    monkeypatch.setattr(document_service, "_container_memory_limit_bytes", lambda: 2 * 1024**3)
    monkeypatch.setattr(document_service, "_run_tesseract_image", fake_ocr)
    progress = []

    result = extract_documents(
        [pptx_path],
        source_language="en",
        progress_callback=lambda completed, total: progress.append((completed, total)),
    )

    assert state["maximum"] == 2
    assert "LEFT OCR PANEL" in result["text"]
    assert "RIGHT OCR PANEL" in result["text"]
    assert "deliberately long speaker note" in result["text"]
    assert result["documents"][0]["ocr_pages"] == 1
    assert result["documents"][0]["ocr_images"] == 2
    assert result["documents"][0]["native_text_pages"] == 0
    assert progress[-1] == (2, 2)


def test_short_title_and_medium_logo_do_not_consume_pptx_ocr_quota(tmp_path: Path, monkeypatch):
    logo_path = tmp_path / "brand-panel.png"
    Image.new("RGB", (800, 400), "navy").save(logo_path)

    pptx_path = tmp_path / "native-title-with-logo.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(5), Inches(0.6))
    title_box.text_frame.text = "Intro"
    slide.shapes.add_picture(str(logo_path), Inches(0.5), Inches(1.4), width=Inches(4))
    presentation.save(pptx_path)

    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda *_args, **_kwargs: pytest.fail("a logo panel must not trigger OCR"),
    )

    result = extract_documents([pptx_path], source_language="en")

    assert "Intro" in result["text"]
    assert result["documents"][0]["ocr_pages"] == 0
    assert result["documents"][0]["ocr_images"] == 0
    assert result["documents"][0]["native_text_pages"] == 1


def test_pptx_auto_ocr_uses_bounded_representative_language_samples(monkeypatch):
    blobs = [bytes([index]) for index in range(9)]
    state = {"active": 0, "maximum": 0}
    detections: list[int] = []
    lock = threading.Lock()
    gate = threading.Barrier(2, timeout=3)

    def fake_detect(blob):
        index = blob[0]
        with lock:
            detections.append(index)
            state["active"] += 1
            state["maximum"] = max(state["maximum"], state["active"])
        if index in {0, 4}:
            gate.wait()
        time.sleep(0.02)
        with lock:
            state["active"] -= 1
        return f"lang-{index}"

    monkeypatch.setattr(document_service, "OCR_PARALLELISM", 2)
    monkeypatch.setattr(document_service, "_container_memory_limit_bytes", lambda: 2 * 1024**3)
    monkeypatch.setattr(document_service, "_detect_pptx_ocr_language", fake_detect)

    languages = document_service._detect_pptx_ocr_languages(blobs)

    assert state["maximum"] == 2
    assert set(detections) == {0, 4, 8}
    assert languages == {
        0: "lang-0",
        1: "lang-0",
        2: "lang-0",
        3: "lang-4",
        4: "lang-4",
        5: "lang-4",
        6: "lang-4",
        7: "lang-8",
        8: "lang-8",
    }


def test_image_only_pdf_is_read_with_ocr(tmp_path: Path, monkeypatch):
    path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with path.open("wb") as stream:
        writer.write(stream)

    monkeypatch.setattr(
        document_service,
        "_render_pdf_page",
        lambda *_: Image.new("L", (1400, 1800), "white"),
    )
    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda image, source_language="auto": (
            "Enerji korunur ve iş sistemler arasında enerji aktarır.",
            "tur+eng",
        ),
    )

    result = extract_documents([path], source_language="tr")
    assert "Enerji korunur" in result["text"]
    assert result["ocr_used"] is True
    assert result["ocr_pages"] == 1
    assert result["documents"][0]["ocr_language"] == "tur+eng"


def test_pdf_ocr_renderer_produces_a_page_image(tmp_path: Path):
    pdf_path = tmp_path / "render-source.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 760, "OCR render safety test")
    pdf.save()

    image = document_service._render_pdf_page(pdf_path, 0)
    try:
        assert image.width >= 1600
        assert image.height >= 2000
    finally:
        image.close()


@pytest.mark.skipif(not shutil.which(config.OCR_COMMAND), reason="Tesseract is not installed")
def test_real_tesseract_reads_a_scanned_pdf(tmp_path: Path):
    image_path = tmp_path / "scan.png"
    image = Image.new("RGB", (1800, 1200), "white")
    draw = ImageDraw.Draw(image)
    font_name = (
        str(Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts" / "arial.ttf")
        if os.name == "nt"
        else "DejaVuSans.ttf"
    )
    font = ImageFont.truetype(font_name, 58)
    draw.text((120, 220), "ENERGY IS CONSERVED IN A CLOSED SYSTEM", fill="black", font=font)
    draw.text((120, 330), "WORK TRANSFERS ENERGY BETWEEN OBJECTS", fill="black", font=font)
    image.save(image_path)

    pdf_path = tmp_path / "real-scan.pdf"
    pdf = canvas.Canvas(str(pdf_path), pagesize=(612, 792))
    pdf.drawImage(str(image_path), 36, 180, width=540, height=360)
    pdf.save()

    result = extract_documents([pdf_path], source_language="en")
    normalized = result["text"].upper()
    assert "ENERGY IS CONSERVED" in normalized
    assert result["ocr_used"] is True
    assert result["ocr_pages"] == 1


def test_scanned_pdf_preflight_estimates_usage_without_running_ocr(tmp_path: Path, monkeypatch):
    path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with path.open("wb") as stream:
        writer.write(stream)
    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda *_args, **_kwargs: pytest.fail("preflight must not run OCR"),
    )

    result = extract_documents([path], enable_ocr=False, allow_ocr_pending=True)
    assert result["ocr_required"] is True
    assert result["estimated"] is True
    assert result["ocr_pages"] == 1
    assert result["words"] == config.OCR_ESTIMATED_WORDS_PER_PAGE
    assert result["credit_minutes"] >= 1


def test_image_document_uses_the_same_ocr_pipeline(tmp_path: Path, monkeypatch):
    path = tmp_path / "whiteboard.png"
    Image.new("RGB", (1200, 800), "white").save(path)
    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda image, source_language="auto": ("Photosynthesis stores chemical energy.", "eng"),
    )

    result = extract_documents([path], source_language="en")
    assert "Photosynthesis" in result["text"]
    assert result["documents"][0]["type"] == "png"
    assert result["ocr_used"] is True


@pytest.mark.parametrize("suffix", [".jpg", ".jpeg", ".webp", ".tif", ".tiff"])
def test_every_supported_image_variant_uses_ocr(tmp_path: Path, monkeypatch, suffix: str):
    path = tmp_path / f"whiteboard{suffix}"
    Image.new("RGB", (640, 360), "white").save(path)
    monkeypatch.setattr(
        document_service,
        "_run_tesseract_image",
        lambda image, source_language="auto": ("Every supported scan reaches OCR.", "eng"),
    )

    result = extract_documents([path], source_language="en")

    assert result["documents"][0]["type"] == suffix.lstrip(".")
    assert result["ocr_used"] is True
    assert "supported scan" in result["text"]


def test_markdown_document_is_extracted(tmp_path: Path):
    path = tmp_path / "lesson.md"
    path.write_text("# Energy\n\nEnergy is conserved in a closed system.", encoding="utf-8")

    result = extract_documents([path], source_language="en")

    assert result["documents"][0]["type"] == "md"
    assert "Energy is conserved" in result["text"]


def test_ocr_page_limit_is_enforced_before_rendering(tmp_path: Path, monkeypatch):
    path = tmp_path / "long-scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.add_blank_page(width=612, height=792)
    with path.open("wb") as stream:
        writer.write(stream)
    monkeypatch.setattr(document_service, "OCR_MAX_PAGES", 1)

    with pytest.raises(LectureSiftError) as caught:
        extract_documents([path])
    assert caught.value.code == "LS-OCR-02"


def test_pdf_ocr_uses_bounded_parallel_pages_and_preserves_page_order(tmp_path: Path, monkeypatch):
    path = tmp_path / "parallel-scan.pdf"
    writer = PdfWriter()
    for _ in range(4):
        writer.add_blank_page(width=612, height=792)
    with path.open("wb") as stream:
        writer.write(stream)

    state = {"active": 0, "maximum": 0}
    lock = threading.Lock()

    def fake_page(_path, page_index, _source_language):
        with lock:
            state["active"] += 1
            state["maximum"] = max(state["maximum"], state["active"])
        time.sleep(0.03)
        with lock:
            state["active"] -= 1
        return page_index, f"PAGE {page_index + 1}", "eng"

    monkeypatch.setattr(document_service, "OCR_PARALLELISM", 2)
    monkeypatch.setattr(document_service, "_ocr_pdf_page", fake_page)
    progress = []
    result = extract_documents([path], progress_callback=lambda done, total: progress.append((done, total)))

    assert state["maximum"] == 2
    assert [result["text"].index(f"PAGE {index}") for index in range(1, 5)] == sorted(
        result["text"].index(f"PAGE {index}") for index in range(1, 5)
    )
    assert progress[-1] == (4, 4)


def test_ocr_parallelism_is_memory_safe_on_small_workers(monkeypatch):
    monkeypatch.setattr(document_service, "OCR_PARALLELISM", 4)
    monkeypatch.setattr(document_service, "_container_memory_limit_bytes", lambda: 512 * 1024 * 1024)
    assert document_service.effective_ocr_parallelism() == 1

    monkeypatch.setattr(document_service, "_container_memory_limit_bytes", lambda: 2 * 1024 * 1024 * 1024)
    assert document_service.effective_ocr_parallelism() == 4


def test_pdf_auto_ocr_uses_bounded_representative_language_samples(tmp_path: Path, monkeypatch):
    path = tmp_path / "multilingual-scan.pdf"
    writer = PdfWriter()
    for _ in range(9):
        writer.add_blank_page(width=612, height=792)
    with path.open("wb") as stream:
        writer.write(stream)

    detections: list[int] = []
    languages: list[str] = []
    monkeypatch.setattr(document_service, "OCR_PARALLELISM", 1)
    monkeypatch.setattr(
        document_service,
        "_detect_pdf_ocr_language",
        lambda _path, page_index: detections.append(page_index) or {
            0: "eng+tur",
            4: "jpn+eng",
            8: "ara+eng",
        }[page_index],
    )

    def fake_page(_path, page_index, language):
        languages.append(language)
        return page_index, f"PAGE {page_index + 1}", language

    monkeypatch.setattr(document_service, "_ocr_pdf_page", fake_page)
    result = extract_documents([path], source_language="auto")

    assert detections == [0, 4, 8]
    assert languages == [
        "eng+tur", "eng+tur", "eng+tur",
        "jpn+eng", "jpn+eng", "jpn+eng", "jpn+eng",
        "ara+eng", "ara+eng",
    ]
    assert result["documents"][0]["ocr_language"] == "eng+tur"


def test_pdf_language_samples_are_detected_concurrently_with_a_bound(tmp_path: Path, monkeypatch):
    path = tmp_path / "language-samples.pdf"
    path.write_bytes(b"%PDF-placeholder")
    state = {"active": 0, "maximum": 0}
    lock = threading.Lock()
    gate = threading.Barrier(2, timeout=3)

    def fake_detect(_path, page_index):
        with lock:
            state["active"] += 1
            state["maximum"] = max(state["maximum"], state["active"])
        if page_index in {0, 4}:
            gate.wait()
        time.sleep(0.02)
        with lock:
            state["active"] -= 1
        return f"lang-{page_index}"

    monkeypatch.setattr(document_service, "OCR_PARALLELISM", 2)
    monkeypatch.setattr(document_service, "_detect_pdf_ocr_language", fake_detect)

    detected = document_service._detect_pdf_ocr_languages(path, list(range(9)))

    assert state["maximum"] == 2
    assert detected[0] == "lang-0"
    assert detected[4] == "lang-4"
    assert detected[8] == "lang-8"


def test_ocr_page_budget_applies_across_all_uploaded_documents(tmp_path: Path, monkeypatch):
    paths = []
    for index in range(2):
        path = tmp_path / f"scan-{index}.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with path.open("wb") as stream:
            writer.write(stream)
        paths.append(path)

    monkeypatch.setattr(document_service, "OCR_MAX_PAGES", 1)
    monkeypatch.setattr(
        document_service,
        "_ocr_pdf_page",
        lambda _path, page_index, _language: (page_index, "Readable scan text", "eng"),
    )
    with pytest.raises(LectureSiftError) as caught:
        extract_documents(paths)
    assert caught.value.code == "LS-OCR-02"


def test_document_upload_releases_response_before_background_quota_preflight(tmp_path: Path, monkeypatch):
    captured: dict = {}

    class DeferredThread:
        def __init__(self, target, args, daemon):
            captured.update(target=target, args=args, daemon=daemon)

        def start(self):
            captured["started"] = True

    monkeypatch.setattr(app_module, "WORK_DIR", tmp_path)
    monkeypatch.setattr(app_module.threading, "Thread", DeferredThread)
    response = TestClient(app).post(
        "/jobs",
        files={"file": ("lecture.txt", b"energy work force power " * 80, "text/plain")},
        headers=_account_headers(),
    )
    assert response.status_code == 200
    body = response.json()
    job = JOBS.get(body["job_id"])
    assert body["source_layout"] == "documents"
    assert body["document_words"] is None
    assert body["billable_minutes"] is None
    assert job["source_type"] == "document"
    assert captured["started"] is True
    assert captured["args"][2]["document_mode"] is True
    seconds = _preflight_documents(body["job_id"], list(captured["args"][1]), captured["args"][2])
    assert seconds is not None and seconds >= 120
    inspected = JOBS.get(body["job_id"])
    assert inspected["document_words"] >= 300
    assert inspected["billable_minutes"] >= 2
    assert captured["args"][2]["document_credit_seconds"] == seconds
    shutil.rmtree(Path(job["job_dir"]), ignore_errors=True)


def test_scanned_pdf_upload_queues_background_ocr(tmp_path: Path, monkeypatch):
    captured: dict = {}

    class DeferredThread:
        def __init__(self, target, args, daemon):
            captured.update(target=target, args=args, daemon=daemon)

        def start(self):
            captured["started"] = True

    monkeypatch.setattr(app_module, "WORK_DIR", tmp_path)
    monkeypatch.setattr(app_module.threading, "Thread", DeferredThread)
    pdf_path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with pdf_path.open("wb") as stream:
        writer.write(stream)

    with pdf_path.open("rb") as stream:
        response = TestClient(app).post(
            "/jobs",
            files={"files": (pdf_path.name, stream, "application/pdf")},
            headers=_account_headers(),
        )

    assert response.status_code == 200, response.text
    body = response.json()
    assert body["ocr_required"] is False
    assert body["ocr_pages"] == 0
    assert body["usage_estimated"] is True
    assert captured["started"] is True
    seconds = _preflight_documents(body["job_id"], list(captured["args"][1]), captured["args"][2])
    assert seconds is not None
    assert captured["args"][2]["document_ocr_required"] is True
    inspected = JOBS.get(body["job_id"])
    assert inspected["document_ocr_pages"] == 1
    job = JOBS.get(body["job_id"])
    shutil.rmtree(Path(job["job_dir"]), ignore_errors=True)


def test_guest_workspace_accepts_a_text_pdf_without_registered_account(tmp_path: Path, monkeypatch):
    captured: dict = {}
    install_rollout_routes(app)

    class DeferredThread:
        def __init__(self, target, args, daemon):
            captured.update(target=target, args=args, daemon=daemon)

        def start(self):
            captured["started"] = True

    monkeypatch.setattr(app_module.threading, "Thread", DeferredThread)
    pdf_path = tmp_path / "guest-lecture.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 760, "Energy is conserved and work transfers energy between systems.")
    pdf.save()

    client = TestClient(app)
    guest = client.post(
        "/billing/guest-session",
        json={"device_id": f"document-browser-{uuid.uuid4()}"},
    )
    assert guest.status_code == 200, guest.text
    with pdf_path.open("rb") as stream:
        response = client.post(
            "/jobs",
            files={"files": (pdf_path.name, stream, "application/pdf")},
            headers={"Authorization": f"Bearer {guest.json()['token']}"},
        )

    assert response.status_code == 200, response.text
    body = response.json()
    assert body["source_layout"] == "documents"
    assert body["document_words"] is None
    assert body["billable_minutes"] is None
    assert captured["started"] is True
    seconds = _preflight_documents(body["job_id"], list(captured["args"][1]), captured["args"][2])
    assert seconds == 60
    inspected = JOBS.get(body["job_id"])
    assert inspected["document_words"] >= 8
    assert inspected["billable_minutes"] == 1
    job = JOBS.get(body["job_id"])
    shutil.rmtree(Path(job["job_dir"]), ignore_errors=True)


def test_document_and_video_cannot_be_mixed_in_one_job():
    response = TestClient(app).post(
        "/jobs",
        files=[
            ("files", ("lesson.txt", b"lesson text", "text/plain")),
            ("files", ("lesson.mp4", b"video", "video/mp4")),
        ],
        headers=_account_headers(),
    )
    assert response.status_code == 400
    assert response.json()["detail"]["code"] == "LS-UPLOAD-05"


def test_cost_ledger_attributes_provider_usage_without_storing_payload(monkeypatch):
    job_id = f"cost-{uuid.uuid4()}"
    response = SimpleNamespace(
        usage=SimpleNamespace(
            prompt_tokens=1000,
            completion_tokens=500,
            prompt_tokens_details=SimpleNamespace(cached_tokens=200),
        )
    )
    try:
        with costs.cost_context(job_id, str(uuid.uuid4())):
            wrote = costs.record_openai_response("gpt-4o-mini", response, "study_pack")
            costs.record_cost(
                provider="test",
                service="guard",
                resource="metadata",
                quantity=1,
                unit="event",
                price_usd=0,
                pricing_source="test",
                pricing_effective_at="2026-08-29",
                metadata={"request_kind": "study_pack", "prompt": "NEVER STORE THIS"},
            )
        assert wrote is True
        with costs.ENGINE.connect() as connection:
            rows = connection.execute(
                costs.COST_EVENTS.select().where(costs.COST_EVENTS.c.job_id == job_id)
            ).mappings().all()
        assert len(rows) == 4
        assert sum(int(row["cost_microusd"]) for row in rows) == 435
        assert all("NEVER STORE THIS" not in row["metadata_json"] for row in rows)
        assert all("prompt" not in json.loads(row["metadata_json"]) for row in rows)
    finally:
        costs.init_cost_database()
        with costs.ENGINE.begin() as connection:
            connection.execute(costs.COST_EVENTS.delete().where(costs.COST_EVENTS.c.job_id == job_id))


def test_r2_roundtrip_costs_can_be_bound_to_a_user_without_a_job():
    user_id = str(uuid.uuid4())
    try:
        with costs.cost_context(None, user_id):
            costs.record_r2_operation("write", bytes_count=61)
            costs.record_r2_operation("read", bytes_count=61)

        with costs.ENGINE.connect() as connection:
            rows = connection.execute(
                costs.COST_EVENTS.select().where(
                    costs.COST_EVENTS.c.user_id == user_id
                )
            ).mappings().all()

        assert len(rows) == 2
        assert all(row["job_id"] is None for row in rows)
        assert {row["resource"] for row in rows} == {
            "class_a_operation",
            "class_b_operation",
        }
        assert {json.loads(row["metadata_json"])["operation"] for row in rows} == {
            "write",
            "read",
        }
    finally:
        costs.init_cost_database()
        with costs.ENGINE.begin() as connection:
            connection.execute(
                costs.COST_EVENTS.delete().where(
                    costs.COST_EVENTS.c.user_id == user_id
                )
            )


def test_admin_cost_endpoint_is_protected_and_separates_external_invoices(monkeypatch):
    install_rollout_routes(app)
    monkeypatch.setattr(config, "ADMIN_ADMIN", "cost-admin-secret")
    monkeypatch.setattr(costs, "_fx_rate", lambda: (42.0, "test rate"))
    client = TestClient(app)
    assert client.get("/billing/admin/costs").status_code == 401
    response = client.get(
        "/billing/admin/costs?days=30&limit=10",
        headers={"Authorization": "Bearer cost-admin-secret"},
    )
    assert response.status_code == 200
    body = response.json()
    assert body["currency"]["source"] == "test rate"
    assert {item["provider"] for item in body["external_invoice_sources"]} >= {
        "iyzico / PayTR",
        "Google Ads",
        "Cloudflare R2",
    }
    assert "kesin kaynaktır" in body["disclaimer"]
    assert body["accuracy"]["required_provider_days"] == len(body["accuracy"]["active_providers"]) * 30
    assert "by_resource" in body and "unit_economics" in body and "actual_costs" in body


def test_invoice_cost_records_are_admin_only_idempotent_and_deletable(monkeypatch):
    install_rollout_routes(app)
    monkeypatch.setattr(config, "ADMIN_ADMIN", "actual-cost-secret")
    monkeypatch.setattr(costs, "_fx_rate", lambda: (40.0, "test rate"))
    client = TestClient(app)
    today = datetime.now(timezone.utc).astimezone(costs._BUSINESS_TIMEZONE).date()
    provider = f"test_{uuid.uuid4().hex[:12]}"
    payload = {
        "provider": provider,
        "service": "monthly_invoice",
        "period_start": today.isoformat(),
        "period_end": today.isoformat(),
        "currency": "TRY",
        "subtotal_minor": 10000,
        "tax_minor": 2000,
        "label": "Test faturası",
        "source_reference": f"INV-{uuid.uuid4().hex[:8]}",
    }
    assert client.post("/billing/admin/costs/actuals", json=payload).status_code == 401
    headers = {"Authorization": "Bearer actual-cost-secret"}
    created = client.post("/billing/admin/costs/actuals", json=payload, headers=headers)
    assert created.status_code == 200, created.text
    actual_id = created.json()["id"]
    try:
        payload["tax_minor"] = 2500
        updated = client.post("/billing/admin/costs/actuals", json=payload, headers=headers)
        assert updated.status_code == 200
        assert updated.json() == {
            "ok": True,
            "id": actual_id,
            "updated": True,
            "message": "Fatura/mutabakat gideri kaydedildi.",
        }
        overview = client.get("/billing/admin/costs?days=1", headers=headers).json()
        saved = next(item for item in overview["actual_costs"] if item["id"] == actual_id)
        assert saved["total_minor"] == 12500
        assert saved["total_usd"] == 3.125
        assert {item["currency"]: item["total_minor"] for item in overview["actual_by_currency"]}["TRY"] >= 12500
    finally:
        deleted = client.delete(f"/billing/admin/costs/actuals/{actual_id}", headers=headers)
        assert deleted.status_code == 200
    assert client.delete(f"/billing/admin/costs/actuals/{actual_id}", headers=headers).status_code == 404


def test_invoice_coverage_merges_overlaps_without_double_counting():
    start = date(2026, 8, 1)
    end = date(2026, 8, 10)
    rows = [
        {"provider": "render", "period_start": "2026-07-25", "period_end": "2026-08-03"},
        {"provider": "render", "period_start": "2026-08-03", "period_end": "2026-08-07"},
        {"provider": "render", "period_start": "2026-08-09", "period_end": "2026-08-20"},
        {"provider": "netlify", "period_start": "2026-08-01", "period_end": "2026-08-10"},
    ]
    assert costs._covered_days(rows, "render", start, end) == 9
    assert costs._covered_days(rows, "netlify", start, end) == 10
    assert costs._covered_days(rows, "resend", start, end) == 0


def test_cost_report_uses_turkiye_calendar_date_near_utc_midnight(monkeypatch):
    class FixedDatetime(datetime):
        @classmethod
        def now(cls, tz=None):
            value = cls(2026, 8, 29, 21, 30, tzinfo=timezone.utc)
            return value if tz is None else value.astimezone(tz)

    monkeypatch.setattr(costs, "datetime", FixedDatetime)
    monkeypatch.setattr(costs, "_fx_rate", lambda: (40.0, "test rate"))
    overview = costs.cost_overview(days=1, limit=1)
    assert overview["period"]["invoice_start"] == "2026-08-30"
    assert overview["period"]["invoice_end"] == "2026-08-30"


def test_cost_report_normalizes_protected_transfer_provider(monkeypatch):
    monkeypatch.setattr(costs, "_fx_rate", lambda: (40.0, "test rate"))
    created = register_user(
        f"cost-provider-{uuid.uuid4()}@example.com",
        "Strong-test-password1",
        "Cost",
        "Provider",
    )
    order = create_payment_order(
        created["user"]["id"],
        IYZICO_BANK_TRANSFER_PROVIDER,
        "plus",
        "monthly",
        "TRY",
    )
    try:
        with costs.ENGINE.begin() as connection:
            connection.execute(
                PAYMENT_ORDERS.update()
                .where(PAYMENT_ORDERS.c.reference == order["reference"])
                .values(status="paid")
            )
        overview = costs.cost_overview(days=1, limit=1)
        active = overview["accuracy"]["active_providers"]
        assert "iyzico" in active
        assert IYZICO_BANK_TRANSFER_PROVIDER not in active
    finally:
        with costs.ENGINE.begin() as connection:
            connection.execute(
                PAYMENT_ORDERS.delete().where(
                    PAYMENT_ORDERS.c.reference == order["reference"]
                )
            )


def test_fixed_cost_confirmation_keys_match_render_configuration():
    services = {item["provider"]: item for item in costs._fixed_services()}
    assert services["domain"]["confirmation_key"] == "LECTURESIFT_COST_DOMAIN_CONFIRMED"
    assert services["render"]["confirmation_key"] == "LECTURESIFT_COST_RENDER_CONFIRMED"
