"""Safe text extraction for supported study documents."""

from __future__ import annotations

import math
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from io import BytesIO
from pathlib import Path
from typing import Any, Callable

from docx import Document
from PIL import Image, ImageOps, UnidentifiedImageError
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import (
    DOCUMENT_EXTENSIONS,
    DOCUMENT_WORDS_PER_CREDIT_MINUTE,
    MAX_DOCUMENT_BYTES,
    MAX_DOCUMENT_CHARACTERS,
    MAX_DOCUMENT_PAGES,
    OCR_COMMAND,
    OCR_DPI,
    OCR_ENABLED,
    OCR_ESTIMATED_WORDS_PER_PAGE,
    OCR_MAX_PAGES,
    OCR_MIN_NATIVE_CHARACTERS,
    OCR_PARALLELISM,
    OCR_PAGE_TIMEOUT_SECONDS,
)
from .errors import LectureSiftError


_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}
_OCR_LANGUAGES = {
    "tr": "tur+eng",
    "en": "eng",
    "de": "deu+eng",
    "fr": "fra+eng",
    "es": "spa+eng",
    "it": "ita+eng",
    "pt": "por+eng",
    "ru": "rus+eng",
    "ar": "ara+eng",
    "zh": "chi_sim+eng",
    "ja": "jpn+eng",
    "ko": "kor+eng",
    "hi": "hin+eng",
}
_OCR_SCRIPT_LANGUAGES = {
    "Arabic": "ara+eng",
    "Cyrillic": "rus+eng",
    "Devanagari": "hin+eng",
    "Han": "chi_sim+eng",
    "Hangul": "kor+eng",
    "Japanese": "jpn+eng",
    "Latin": "eng+tur+deu+fra+spa+ita+por",
}
_MAX_IMAGE_PIXELS = 40_000_000
_ONE_GIB = 1024 * 1024 * 1024
_PPTX_MIN_PANEL_COVERAGE = 0.05
_PPTX_MIN_SCANNED_SLIDE_COVERAGE = 0.30


def _container_memory_limit_bytes() -> int | None:
    """Return the container memory limit without depending on psutil.

    Render uses cgroup v2 today, while older/local Linux installations may
    still expose the v1 path.  An unbounded value is treated as unknown.
    """
    candidates = (
        Path("/sys/fs/cgroup/memory.max"),
        Path("/sys/fs/cgroup/memory/memory.limit_in_bytes"),
    )
    for candidate in candidates:
        try:
            raw = candidate.read_text(encoding="ascii").strip()
            if raw == "max":
                return None
            value = int(raw)
            if 0 < value < 1 << 60:
                return value
        except (OSError, ValueError):
            continue
    return None


def effective_ocr_parallelism() -> int:
    """Keep OCR parallel on capable workers and memory-safe on small ones."""
    memory_limit = _container_memory_limit_bytes()
    if memory_limit is not None and memory_limit < _ONE_GIB:
        return 1
    return OCR_PARALLELISM


def _normalize_text(value: str) -> str:
    value = value.replace("\x00", "")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in value.splitlines()]
    output: list[str] = []
    blank = False
    for line in lines:
        if line:
            output.append(line)
            blank = False
        elif output and not blank:
            output.append("")
            blank = True
    return "\n".join(output).strip()


def _validate_size(path: Path) -> None:
    size = path.stat().st_size
    if size <= 0:
        raise LectureSiftError("LS-DOC-01", "Belge boş görünüyor.", f"Empty document: {path.name}", 400)
    if size > MAX_DOCUMENT_BYTES:
        raise LectureSiftError(
            "LS-DOC-02",
            f"Bir belge en fazla {MAX_DOCUMENT_BYTES // (1024 * 1024)} MB olabilir.",
            f"Document exceeds size limit: {path.name}",
            413,
        )


def _validate_office_archive(path: Path) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
            if len(members) > 10_000:
                raise ValueError("too many archive members")
            total_uncompressed = 0
            for member in members:
                name = member.filename.replace("\\", "/")
                if name.startswith("/") or ".." in name.split("/"):
                    raise ValueError("unsafe archive path")
                total_uncompressed += max(0, member.file_size)
                if member.compress_size and member.file_size / member.compress_size > 250:
                    raise ValueError("unsafe compression ratio")
            if total_uncompressed > 200 * 1024 * 1024:
                raise ValueError("archive expansion limit")
    except (zipfile.BadZipFile, ValueError) as exc:
        raise LectureSiftError(
            "LS-DOC-03",
            "Word veya PowerPoint belgesi bozuk ya da güvenli açılamıyor.",
            f"Unsafe Office archive {path.name}: {exc}",
            400,
        ) from exc


def _require_ocr() -> None:
    if not OCR_ENABLED:
        raise LectureSiftError(
            "LS-OCR-01",
            "Bu sunucuda OCR geçici olarak kapalı. Biraz sonra yeniden deneyebilirsin.",
            "OCR is disabled by configuration",
            503,
        )
    if not shutil.which(OCR_COMMAND):
        raise LectureSiftError(
            "LS-OCR-01",
            "OCR hizmeti şu anda kullanılamıyor. Biraz sonra yeniden deneyebilirsin.",
            "Tesseract executable is unavailable",
            503,
        )


def _validate_image_dimensions(image: Image.Image) -> None:
    if image.width <= 0 or image.height <= 0 or image.width * image.height > _MAX_IMAGE_PIXELS:
        raise LectureSiftError(
            "LS-OCR-04",
            "Görsel güvenli OCR sınırlarını aşıyor. Daha düşük çözünürlüklü bir kopya yükle.",
            f"Unsafe OCR image dimensions: {image.width}x{image.height}",
            413,
        )


def _safe_image(image: Image.Image) -> Image.Image:
    _validate_image_dimensions(image)
    image.load()
    grayscale = ImageOps.autocontrast(ImageOps.grayscale(image))
    if grayscale.width < 1400:
        ratio = min(2.0, 1400 / max(1, grayscale.width))
        grayscale = grayscale.resize(
            (round(grayscale.width * ratio), round(grayscale.height * ratio)),
            Image.Resampling.LANCZOS,
        )
    return grayscale


def _auto_ocr_languages(image_path: Path) -> str:
    try:
        completed = subprocess.run(
            [OCR_COMMAND, str(image_path), "stdout", "-l", "osd", "--psm", "0"],
            capture_output=True,
            text=True,
            timeout=min(20, OCR_PAGE_TIMEOUT_SECONDS),
            check=False,
        )
        match = re.search(r"^Script:\s*(.+?)\s*$", completed.stdout, flags=re.MULTILINE)
        if match:
            return _OCR_SCRIPT_LANGUAGES.get(match.group(1).strip(), "eng+tur")
    except (OSError, subprocess.TimeoutExpired):
        pass
    return "eng+tur"


def _run_tesseract_image(image: Image.Image, source_language: str = "auto") -> tuple[str, str]:
    _require_ocr()
    prepared = _safe_image(image)
    try:
        with tempfile.TemporaryDirectory(prefix="lecturesift-ocr-") as directory:
            # Tesseract accepts portable graymap directly.  It avoids the
            # relatively expensive PNG compression pass for every page and
            # keeps peak memory lower on small workers.
            image_path = Path(directory) / "page.pgm"
            prepared.save(image_path, format="PPM")
            # PDF extraction can pass a Tesseract language expression that was
            # detected once for the whole document. This avoids one relatively
            # expensive OSD process for every scanned page.
            language = _OCR_LANGUAGES.get(source_language)
            if not language:
                language = _auto_ocr_languages(image_path) if source_language == "auto" else source_language
            try:
                completed = subprocess.run(
                    [
                        OCR_COMMAND,
                        str(image_path),
                        "stdout",
                        "-l",
                        language,
                        "--psm",
                        "3",
                        "--dpi",
                        str(OCR_DPI),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=OCR_PAGE_TIMEOUT_SECONDS,
                    check=False,
                    env={**os.environ, "OMP_THREAD_LIMIT": "1"},
                )
            except subprocess.TimeoutExpired as exc:
                raise LectureSiftError(
                    "LS-OCR-03",
                    "Bu sayfanın OCR işlemi zaman sınırını aştı. Belgeyi bölerek yeniden dene.",
                    "Tesseract page timeout",
                    422,
                ) from exc
            except OSError as exc:
                raise LectureSiftError(
                    "LS-OCR-01",
                    "OCR hizmeti şu anda kullanılamıyor. Biraz sonra yeniden deneyebilirsin.",
                    f"Tesseract launch error: {type(exc).__name__}",
                    503,
                ) from exc
            if completed.returncode != 0:
                raise LectureSiftError(
                    "LS-OCR-01",
                    "OCR motoru belge dilini işleyemedi. Kaynak dilini seçip yeniden dene.",
                    f"Tesseract exited with code {completed.returncode}",
                    503,
                )
    finally:
        prepared.close()
    return _normalize_text(completed.stdout), language


def _detect_pdf_ocr_language(path: Path, page_index: int) -> str:
    """Detect a scan's script once before parallel page OCR begins."""
    image = _render_pdf_page(path, page_index)
    prepared: Image.Image | None = None
    try:
        prepared = _safe_image(image)
        with tempfile.TemporaryDirectory(prefix="lecturesift-ocr-language-") as directory:
            image_path = Path(directory) / "sample.pgm"
            prepared.save(image_path, format="PPM")
            return _auto_ocr_languages(image_path)
    finally:
        if prepared is not None:
            prepared.close()
        image.close()


def _detect_pdf_ocr_languages(path: Path, page_indexes: list[int]) -> dict[int, str]:
    """Map pages to at most three representative script detections.

    Long scanned documents may change language between chapters. Sampling the
    first, middle, and last OCR page preserves those common transitions while
    replacing up to hundreds of per-page OSD subprocesses with at most three.
    """
    if not page_indexes:
        return {}
    positions = sorted({0, len(page_indexes) // 2, len(page_indexes) - 1})
    sample_pages = [page_indexes[position] for position in positions]
    workers = min(effective_ocr_parallelism(), len(sample_pages))
    if workers <= 1:
        samples = {
            page_index: _detect_pdf_ocr_language(path, page_index)
            for page_index in sample_pages
        }
    else:
        with ThreadPoolExecutor(max_workers=workers, thread_name_prefix="lecturesift-ocr-language") as executor:
            futures = {
                executor.submit(_detect_pdf_ocr_language, path, page_index): page_index
                for page_index in sample_pages
            }
            samples = {futures[future]: future.result() for future in as_completed(futures)}
    return {
        page_index: samples[min(samples, key=lambda sampled: abs(sampled - page_index))]
        for page_index in page_indexes
    }


def _render_pdf_page(path: Path, page_index: int) -> Image.Image:
    try:
        import pypdfium2 as pdfium

        document = pdfium.PdfDocument(str(path))
        try:
            page = document[page_index]
            try:
                bitmap = page.render(scale=OCR_DPI / 72)
                try:
                    return bitmap.to_pil().copy()
                finally:
                    bitmap.close()
            finally:
                page.close()
        finally:
            document.close()
    except LectureSiftError:
        raise
    except Exception as exc:
        raise LectureSiftError(
            "LS-OCR-04",
            "PDF sayfası OCR için görüntüye dönüştürülemedi.",
            f"PDF render failed on page {page_index + 1}: {type(exc).__name__}",
            422,
        ) from exc


def _ocr_pdf_page(path: Path, page_index: int, source_language: str) -> tuple[int, str, str]:
    """Render and OCR one page without sharing mutable PDF/image state."""
    image = _render_pdf_page(path, page_index)
    try:
        text, language = _run_tesseract_image(image, source_language)
    finally:
        image.close()
    return page_index, text, language


def _extract_pdf_native_pages(path: Path, page_count: int) -> tuple[list[str], str]:
    """Use PDFium's native text engine, with pypdf as a compatibility fallback."""
    try:
        import pypdfium2 as pdfium

        document = pdfium.PdfDocument(str(path))
        try:
            if len(document) != page_count:
                raise ValueError("PDF page count changed between readers")
            parts: list[str] = []
            for page_index in range(page_count):
                page = document[page_index]
                try:
                    text_page = page.get_textpage()
                    try:
                        parts.append(_normalize_text(text_page.get_text_range() or ""))
                    finally:
                        text_page.close()
                finally:
                    page.close()
            return parts, "pdfium"
        finally:
            document.close()
    except Exception:
        reader = PdfReader(str(path), strict=False)
        return [_normalize_text(page.extract_text() or "") for page in reader.pages], "pypdf"


def _extract_pdf(
    path: Path,
    source_language: str = "auto",
    *,
    enable_ocr: bool = True,
    progress_callback: Callable[[int, int], None] | None = None,
    ocr_page_budget: int | None = None,
) -> tuple[str, dict[str, Any]]:
    with path.open("rb") as stream:
        signature = stream.read(5)
    if signature != b"%PDF-":
        raise LectureSiftError("LS-DOC-04", "PDF dosyası geçerli görünmüyor.", f"Invalid PDF signature: {path.name}", 400)
    try:
        reader = PdfReader(str(path), strict=False)
        if reader.is_encrypted:
            raise LectureSiftError("LS-DOC-05", "Şifreli PDF dosyaları desteklenmiyor.", "Encrypted PDF", 400)
        if len(reader.pages) > MAX_DOCUMENT_PAGES:
            raise LectureSiftError(
                "LS-DOC-06",
                f"PDF en fazla {MAX_DOCUMENT_PAGES} sayfa olabilir.",
                f"PDF page limit exceeded: {len(reader.pages)}",
                413,
            )
        page_count = len(reader.pages)
        parts, text_engine = _extract_pdf_native_pages(path, page_count)
    except LectureSiftError:
        raise
    except Exception as exc:
        raise LectureSiftError("LS-DOC-04", "PDF güvenli biçimde okunamadı.", str(exc), 400) from exc
    ocr_indexes = [index for index, part in enumerate(parts) if len(part) < OCR_MIN_NATIVE_CHARACTERS]
    native_text_pages = len(parts) - len(ocr_indexes)
    ocr_language = ""
    available_ocr_pages = OCR_MAX_PAGES if ocr_page_budget is None else max(0, min(OCR_MAX_PAGES, ocr_page_budget))
    if len(ocr_indexes) > available_ocr_pages:
        raise LectureSiftError(
            "LS-OCR-02",
            f"Tek işte toplam en fazla {OCR_MAX_PAGES} taranmış sayfaya OCR uygulanabilir. Belgeyi bölerek yükle.",
            f"OCR page limit exceeded: {len(ocr_indexes)} pages with {available_ocr_pages} remaining",
            413,
        )
    if enable_ocr and ocr_indexes:
        results: dict[int, tuple[str, str]] = {}
        ocr_languages = (
            _detect_pdf_ocr_languages(path, ocr_indexes)
            if source_language == "auto" and len(ocr_indexes) > 1
            else {page_index: source_language for page_index in ocr_indexes}
        )
        workers = min(effective_ocr_parallelism(), len(ocr_indexes))
        if workers == 1:
            completed_results = (
                _ocr_pdf_page(path, page_index, ocr_languages[page_index]) for page_index in ocr_indexes
            )
            for completed_pages, (page_index, text, language) in enumerate(completed_results, 1):
                results[page_index] = (text, language)
                if progress_callback:
                    progress_callback(completed_pages, len(ocr_indexes))
        else:
            with ThreadPoolExecutor(max_workers=workers, thread_name_prefix="lecturesift-ocr") as executor:
                futures = {
                    executor.submit(_ocr_pdf_page, path, page_index, ocr_languages[page_index]): page_index
                    for page_index in ocr_indexes
                }
                for completed_pages, future in enumerate(as_completed(futures), 1):
                    page_index, text, language = future.result()
                    results[page_index] = (text, language)
                    if progress_callback:
                        progress_callback(completed_pages, len(ocr_indexes))
        for page_index in ocr_indexes:
            text, language = results[page_index]
            if len(text) > len(parts[page_index]):
                parts[page_index] = text
            if not ocr_language and language:
                ocr_language = language
    return "\n\n".join(part for part in parts if part), {
        "pages": len(reader.pages),
        "native_text_pages": native_text_pages,
        "ocr_pages": len(ocr_indexes),
        "ocr_used": bool(enable_ocr and ocr_indexes),
        "ocr_required": bool(not enable_ocr and ocr_indexes),
        "ocr_language": ocr_language,
        "text_engine": text_engine,
    }


def _extract_image(
    path: Path,
    source_language: str = "auto",
    *,
    enable_ocr: bool = True,
    progress_callback: Callable[[int, int], None] | None = None,
    ocr_page_budget: int | None = None,
) -> tuple[str, dict[str, Any]]:
    if ocr_page_budget is not None and ocr_page_budget < 1:
        raise LectureSiftError(
            "LS-OCR-02",
            f"Tek işte toplam en fazla {OCR_MAX_PAGES} taranmış sayfaya OCR uygulanabilir. Dosyaları bölerek yükle.",
            "OCR image exceeds remaining page budget",
            413,
        )
    try:
        with Image.open(path) as source:
            _validate_image_dimensions(source)
            if not enable_ocr:
                return "", {
                    "pages": 1,
                    "native_text_pages": 0,
                    "ocr_pages": 1,
                    "ocr_used": False,
                    "ocr_required": True,
                    "ocr_language": "",
                }
            image = source.copy()
    except (Image.DecompressionBombError, UnidentifiedImageError, OSError) as exc:
        raise LectureSiftError(
            "LS-OCR-04",
            "Görsel dosyası güvenli biçimde okunamadı.",
            f"Unreadable image: {type(exc).__name__}",
            422,
        ) from exc
    try:
        text, language = _run_tesseract_image(image, source_language)
    finally:
        image.close()
    if progress_callback:
        progress_callback(1, 1)
    return text, {
        "pages": 1,
        "native_text_pages": 0,
        "ocr_pages": 1,
        "ocr_used": True,
        "ocr_required": False,
        "ocr_language": language,
    }


def _extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
    _validate_office_archive(path)
    try:
        document = Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
    except Exception as exc:
        raise LectureSiftError("LS-DOC-07", "Word belgesi okunamadı.", str(exc), 400) from exc
    return _normalize_text("\n\n".join(parts)), {"paragraphs": len(document.paragraphs), "tables": len(document.tables)}


def _pptx_shape_text(shape: Any) -> list[str]:
    """Return user-authored text from a PowerPoint shape, including tables."""
    parts: list[str] = []
    text = getattr(shape, "text", "")
    if text and text.strip():
        parts.append(text.strip())
    if bool(getattr(shape, "has_table", False)):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            parts.extend(_pptx_shape_text(child))
    return parts


def _pptx_picture_shapes(shape: Any) -> list[Any]:
    # Inserted picture placeholders are ``PLACEHOLDER`` shapes in
    # python-pptx rather than ``PICTURE`` shapes. Capability detection keeps
    # those images visible to OCR without guessing from placeholder types.
    try:
        image = shape.image
        blob = image.blob
    except (AttributeError, KeyError, TypeError, ValueError):
        blob = None
    if isinstance(blob, (bytes, bytearray)) and blob:
        return [shape]
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        pictures: list[Any] = []
        for child in shape.shapes:
            pictures.extend(_pptx_picture_shapes(child))
        return pictures
    return []


def _pptx_picture_blob(shape: Any) -> bytes:
    try:
        blob = shape.image.blob
    except (AttributeError, KeyError, TypeError, ValueError) as exc:
        raise LectureSiftError(
            "LS-OCR-04",
            "Sunumdaki taranmış slayt görseli güvenli biçimde okunamadı.",
            f"Unreadable PowerPoint image relationship: {type(exc).__name__}",
            422,
        ) from exc
    if not isinstance(blob, (bytes, bytearray)) or not blob:
        raise LectureSiftError(
            "LS-OCR-04",
            "Sunumdaki taranmış slayt görseli güvenli biçimde okunamadı.",
            "Empty PowerPoint image payload",
            422,
        )
    return bytes(blob)


def _load_pptx_picture_image(blob: bytes) -> Image.Image:
    try:
        with Image.open(BytesIO(blob)) as source:
            _validate_image_dimensions(source)
            return source.copy()
    except (Image.DecompressionBombError, UnidentifiedImageError, OSError) as exc:
        raise LectureSiftError(
            "LS-OCR-04",
            "Sunumdaki taranmış slayt görseli güvenli biçimde okunamadı.",
            f"Unreadable PowerPoint image: {type(exc).__name__}",
            422,
        ) from exc


def _detect_pptx_ocr_language(blob: bytes) -> str:
    image = _load_pptx_picture_image(blob)
    prepared: Image.Image | None = None
    try:
        prepared = _safe_image(image)
        with tempfile.TemporaryDirectory(prefix="lecturesift-pptx-ocr-language-") as directory:
            image_path = Path(directory) / "sample.pgm"
            prepared.save(image_path, format="PPM")
            return _auto_ocr_languages(image_path)
    finally:
        if prepared is not None:
            prepared.close()
        image.close()


def _detect_pptx_ocr_languages(blobs: list[bytes]) -> dict[int, str]:
    """Map embedded images to at most three representative script samples."""
    if not blobs:
        return {}
    positions = sorted({0, len(blobs) // 2, len(blobs) - 1})
    workers = min(effective_ocr_parallelism(), len(positions))
    if workers <= 1:
        samples = {position: _detect_pptx_ocr_language(blobs[position]) for position in positions}
    else:
        with ThreadPoolExecutor(
            max_workers=workers,
            thread_name_prefix="lecturesift-pptx-ocr-language",
        ) as executor:
            futures = {
                executor.submit(_detect_pptx_ocr_language, blobs[position]): position
                for position in positions
            }
            samples = {futures[future]: future.result() for future in as_completed(futures)}
    return {
        position: samples[
            min(samples, key=lambda sampled: (abs(sampled - position), sampled))
        ]
        for position in range(len(blobs))
    }


def _ocr_pptx_picture(candidate_index: int, blob: bytes, source_language: str) -> tuple[int, str, str]:
    image = _load_pptx_picture_image(blob)
    try:
        text, language = _run_tesseract_image(image, source_language)
    finally:
        image.close()
    return candidate_index, text, language


def _extract_pptx(
    path: Path,
    source_language: str = "auto",
    *,
    enable_ocr: bool = True,
    progress_callback: Callable[[int, int], None] | None = None,
    ocr_page_budget: int | None = None,
) -> tuple[str, dict[str, Any]]:
    _validate_office_archive(path)
    try:
        presentation = Presentation(str(path))
        if len(presentation.slides) > MAX_DOCUMENT_PAGES:
            raise LectureSiftError(
                "LS-DOC-06",
                f"Sunum en fazla {MAX_DOCUMENT_PAGES} slayt olabilir.",
                f"Presentation slide limit exceeded: {len(presentation.slides)}",
                413,
            )
        slide_lines: dict[int, list[str]] = {}
        ocr_candidates: list[tuple[int, bytes]] = []
        slide_area = max(1, int(presentation.slide_width) * int(presentation.slide_height))
        for index, slide in enumerate(presentation.slides, 1):
            visual_lines: list[str] = []
            notes_lines: list[str] = []
            pictures: list[Any] = []
            for shape in slide.shapes:
                visual_lines.extend(_pptx_shape_text(shape))
                pictures.extend(_pptx_picture_shapes(shape))
            if bool(getattr(slide, "has_notes_slide", False)):
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text if notes_frame is not None else ""
                if notes_text and notes_text.strip():
                    notes_lines.append(f"SPEAKER NOTES\n{notes_text.strip()}")
            slide_lines[index] = [*visual_lines, *notes_lines]
            # Notes are invisible during the presentation and must not make a
            # scanned/photo slide look like it already has native slide text.
            if pictures and len(" ".join(visual_lines).strip()) < OCR_MIN_NATIVE_CHARACTERS:
                picture_coverage = [
                    (
                        picture,
                        (
                            max(0, int(getattr(picture, "width", 0)))
                            * max(0, int(getattr(picture, "height", 0)))
                        )
                        / slide_area,
                    )
                    for picture in pictures
                ]
                substantial = [
                    picture
                    for picture, coverage in picture_coverage
                    if coverage >= _PPTX_MIN_PANEL_COVERAGE
                ]
                # A short title plus a logo/photo is still a native slide and
                # must not consume the user's OCR allowance. Treat the slide
                # as scanned only when its meaningful image panels together
                # occupy a material part of the canvas. Multiple screenshot
                # panels are summed so two-column scans remain supported.
                if sum(
                    coverage
                    for _picture, coverage in picture_coverage
                    if coverage >= _PPTX_MIN_PANEL_COVERAGE
                ) < _PPTX_MIN_SCANNED_SLIDE_COVERAGE:
                    substantial = []
                # OCR every meaningful panel on low-text slides. This covers
                # two-column screenshots and picture placeholders while
                # ignoring small logos and decorative icons.
                ocr_candidates.extend(
                    (index, _pptx_picture_blob(picture)) for picture in substantial
                )

        available_ocr_pages = (
            OCR_MAX_PAGES
            if ocr_page_budget is None
            else max(0, min(OCR_MAX_PAGES, ocr_page_budget))
        )
        if len(ocr_candidates) > available_ocr_pages:
            raise LectureSiftError(
                "LS-OCR-02",
                f"Tek işte toplam en fazla {OCR_MAX_PAGES} taranmış sayfaya OCR uygulanabilir. Sunumu bölerek yükle.",
                f"PowerPoint OCR image limit exceeded: {len(ocr_candidates)}",
                413,
            )

        ocr_used = False
        ocr_languages: list[str] = []
        if enable_ocr:
            blobs = [blob for _, blob in ocr_candidates]
            candidate_languages = (
                _detect_pptx_ocr_languages(blobs)
                if source_language == "auto" and len(blobs) > 1
                else {candidate_index: source_language for candidate_index in range(len(blobs))}
            )
            results: dict[int, tuple[str, str]] = {}
            workers = min(effective_ocr_parallelism(), len(ocr_candidates)) if ocr_candidates else 0
            if workers == 1:
                completed_results = (
                    _ocr_pptx_picture(
                        candidate_index,
                        blob,
                        candidate_languages[candidate_index],
                    )
                    for candidate_index, (_, blob) in enumerate(ocr_candidates)
                )
                for completed, (candidate_index, text, language) in enumerate(completed_results, 1):
                    results[candidate_index] = (text, language)
                    if progress_callback:
                        progress_callback(completed, len(ocr_candidates))
            elif workers > 1:
                with ThreadPoolExecutor(
                    max_workers=workers,
                    thread_name_prefix="lecturesift-pptx-ocr",
                ) as executor:
                    futures = {
                        executor.submit(
                            _ocr_pptx_picture,
                            candidate_index,
                            blob,
                            candidate_languages[candidate_index],
                        ): candidate_index
                        for candidate_index, (_, blob) in enumerate(ocr_candidates)
                    }
                    for completed, future in enumerate(as_completed(futures), 1):
                        candidate_index, text, language = future.result()
                        results[candidate_index] = (text, language)
                        if progress_callback:
                            progress_callback(completed, len(ocr_candidates))
            for candidate_index, (index, _) in enumerate(ocr_candidates):
                text, language = results[candidate_index]
                if text.strip():
                    slide_lines[index].append(f"SLIDE OCR\n{text.strip()}")
                if language and language not in ocr_languages:
                    ocr_languages.append(language)
            ocr_used = bool(ocr_candidates)

        sections = [
            f"SLIDE {index}\n" + "\n".join(lines)
            for index, lines in slide_lines.items()
            if lines
        ]
    except LectureSiftError:
        raise
    except Exception as exc:
        raise LectureSiftError("LS-DOC-08", "PowerPoint sunumu okunamadı.", str(exc), 400) from exc
    ocr_slide_indexes = {index for index, _ in ocr_candidates}
    return _normalize_text("\n\n".join(sections)), {
        "slides": len(presentation.slides),
        "native_text_pages": len(presentation.slides) - len(ocr_slide_indexes),
        "ocr_pages": len(ocr_slide_indexes),
        "ocr_images": len(ocr_candidates),
        "ocr_units": len(ocr_candidates),
        "ocr_used": ocr_used,
        "ocr_required": bool(ocr_candidates) and not enable_ocr,
        "ocr_language": "+".join(ocr_languages),
    }


def _extract_text(path: Path) -> tuple[str, dict[str, Any]]:
    raw = path.read_bytes()
    if b"\x00" in raw[:4096]:
        raise LectureSiftError("LS-DOC-09", "Metin dosyası geçerli görünmüyor.", "Binary content in text file", 400)
    for encoding in ("utf-8-sig", "utf-16", "cp1254", "latin-1"):
        try:
            return _normalize_text(raw.decode(encoding)), {"encoding": encoding}
        except UnicodeDecodeError:
            continue
    raise LectureSiftError("LS-DOC-09", "Metin dosyasının karakter kodlaması okunamadı.", "Unknown text encoding", 400)


def extract_documents(
    paths: list[Path],
    source_language: str = "auto",
    *,
    enable_ocr: bool = True,
    allow_ocr_pending: bool = False,
    progress_callback: Callable[[int, int], None] | None = None,
) -> dict[str, Any]:
    if not paths:
        raise LectureSiftError("LS-DOC-10", "En az bir belge ekle.", "No document paths", 400)
    sections: list[str] = []
    metadata: list[dict[str, Any]] = []
    completed_ocr_units = 0
    for index, path in enumerate(paths, 1):
        _validate_size(path)
        suffix = path.suffix.casefold()
        if suffix not in DOCUMENT_EXTENSIONS:
            raise LectureSiftError("LS-DOC-11", "Bu belge biçimi desteklenmiyor.", f"Unsupported document: {suffix}", 400)
        if suffix == ".pdf":
            progress_offset = completed_ocr_units
            text, details = _extract_pdf(
                path,
                source_language,
                enable_ocr=enable_ocr,
                progress_callback=(
                    (lambda completed, total, offset=progress_offset: progress_callback(offset + completed, offset + total))
                    if progress_callback else None
                ),
                ocr_page_budget=OCR_MAX_PAGES - completed_ocr_units,
            )
        elif suffix == ".docx":
            text, details = _extract_docx(path)
        elif suffix == ".pptx":
            progress_offset = completed_ocr_units
            text, details = _extract_pptx(
                path,
                source_language,
                enable_ocr=enable_ocr,
                progress_callback=(
                    (lambda completed, total, offset=progress_offset: progress_callback(offset + completed, offset + total))
                    if progress_callback else None
                ),
                ocr_page_budget=OCR_MAX_PAGES - completed_ocr_units,
            )
        elif suffix in _IMAGE_EXTENSIONS:
            progress_offset = completed_ocr_units
            text, details = _extract_image(
                path,
                source_language,
                enable_ocr=enable_ocr,
                progress_callback=(
                    (lambda completed, total, offset=progress_offset: progress_callback(offset + completed, offset + total))
                    if progress_callback else None
                ),
                ocr_page_budget=OCR_MAX_PAGES - completed_ocr_units,
            )
        else:
            text, details = _extract_text(path)
        ocr_pending = bool(details.get("ocr_required"))
        if not text.strip() and not (allow_ocr_pending and ocr_pending):
            raise LectureSiftError(
                "LS-DOC-12",
                "OCR tamamlandı ancak okunabilir metin bulunamadı. Daha net bir tarama veya doğru kaynak diliyle yeniden dene.",
                f"No extractable text: {path.name}",
                422,
            )
        if text.strip():
            sections.append(f"DOCUMENT {index}: {path.name}\n{text}")
        metadata.append({"name": path.name, "type": suffix.lstrip("."), **details, "characters": len(text)})
        completed_ocr_units += int(details.get("ocr_units") or details.get("ocr_pages") or 0)
    combined = "\n\n".join(sections)
    if len(combined) > MAX_DOCUMENT_CHARACTERS:
        raise LectureSiftError(
            "LS-DOC-13",
            "Belgelerin toplam metni güvenli işleme sınırını aşıyor. Kaynağı bölerek yeniden yükle.",
            f"Document character limit exceeded: {len(combined)}",
            413,
        )
    extracted_words = len(re.findall(r"[^\W_]+", combined, flags=re.UNICODE))
    pending_ocr_pages = sum(
        int(item.get("ocr_pages") or 0) for item in metadata if item.get("ocr_required")
    )
    estimated_ocr_words = pending_ocr_pages * OCR_ESTIMATED_WORDS_PER_PAGE
    words = extracted_words + estimated_ocr_words
    credit_minutes = max(1, math.ceil(words / DOCUMENT_WORDS_PER_CREDIT_MINUTE))
    return {
        "text": combined,
        "documents": metadata,
        "characters": len(combined),
        "words": words,
        "extracted_words": extracted_words,
        "estimated": bool(pending_ocr_pages),
        "ocr_required": bool(pending_ocr_pages),
        "pages": sum(int(item.get("pages") or item.get("slides") or 1) for item in metadata),
        "ocr_pages": sum(int(item.get("ocr_pages") or 0) for item in metadata),
        "ocr_used": any(bool(item.get("ocr_used")) for item in metadata),
        "credit_minutes": credit_minutes,
        "credit_seconds": credit_minutes * 60,
    }
